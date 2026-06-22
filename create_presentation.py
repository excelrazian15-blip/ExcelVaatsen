#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Script untuk membuat presentasi PowerPoint tentang Data Mining
Gunakan: python create_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle):
    """Membuat slide judul"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(41, 128, 185)  # Biru
    
    # Judul
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(236, 240, 241)
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_points):
    """Membuat slide konten dengan bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(41, 128, 185)
    title_shape.line.color.rgb = RGBColor(41, 128, 185)
    
    # Title text
    title_frame = title_shape.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_before = Pt(10)
    p.space_after = Pt(10)
    
    # Content
    text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(44, 62, 80)
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.alignment = PP_ALIGN.LEFT

def create_presentation():
    """Membuat presentasi Data Mining"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    add_title_slide(prs, "DATA MINING", "Pengertian, Fungsi, Tahapan, dan Penerapan")
    
    # Slide 2: Pengertian Data Mining
    add_content_slide(prs, "1. Pengertian Data Mining", [
        "• Data Mining adalah proses ekstraksi pola, informasi, dan pengetahuan",
        "  dari data dalam jumlah besar.",
        "",
        "• Tujuan: Menemukan pola tersembunyi dan hubungan antar data yang",
        "  dapat digunakan untuk pengambilan keputusan.",
        "",
        "• Data Mining menggabungkan teknik dari Statistik, Machine Learning,",
        "  dan Database Management.",
        "",
        "• Juga dikenal sebagai Knowledge Discovery in Databases (KDD)."
    ])
    
    # Slide 3: Pengenalan Fungsi
    add_content_slide(prs, "2. Pengenalan Fungsi Data Mining", [
        "• Prediksi (Prediction): Memprediksi nilai masa depan berdasarkan data",
        "  historis.",
        "",
        "• Klasifikasi (Classification): Mengelompokkan data ke dalam kategori",
        "  yang telah ditentukan.",
        "",
        "• Pengelompokan (Clustering): Mengelompokkan data serupa tanpa label",
        "  yang telah ditentukan.",
        "",
        "• Asosiasi (Association): Menemukan hubungan antar item dalam dataset."
    ])
    
    # Slide 4: Fungsi Data Mining
    add_content_slide(prs, "3. Fungsi Data Mining", [
        "• Deskripsi (Description): Menjelaskan karakteristik data yang ada.",
        "",
        "• Estimasi (Estimation): Memperkirakan nilai variabel berdasarkan data.",
        "",
        "• Prediksi (Prediction): Meramalkan nilai yang akan datang.",
        "",
        "• Klasifikasi (Classification): Menentukan kategori data.",
        "",
        "• Clustering (Clustering): Pengelompokan data tanpa label awal.",
        "",
        "• Association Rules: Mencari aturan asosiasi antar atribut."
    ])
    
    # Slide 5: Pengenalan Tahapan Data Mining
    add_content_slide(prs, "4. Pengenalan Tahapan Data Mining", [
        "1. Problem Definition: Mendefinisikan masalah dan tujuan",
        "",
        "2. Data Collection: Mengumpulkan data dari berbagai sumber",
        "",
        "3. Data Preparation: Membersihkan dan mempersiapkan data",
        "",
        "4. Exploration & Visualization: Menjelajahi dan visualisasi data",
        "",
        "5. Modeling: Membangun model dengan algoritma tertentu",
        "",
        "6. Evaluation: Mengevaluasi hasil model",
        "",
        "7. Deployment: Implementasi model di lapangan"
    ])
    
    # Slide 6: Penerapan Data Mining
    add_content_slide(prs, "5. Penerapan Data Mining", [
        "• E-Commerce: Rekomendasi produk berdasarkan perilaku pembeli",
        "",
        "• Perbankan: Deteksi fraud dan credit scoring",
        "",
        "• Healthcare: Diagnosis penyakit dan prediksi penyebaran penyakit",
        "",
        "• Marketing: Segmentasi pelanggan dan target marketing",
        "",
        "• Manufaktur: Quality control dan predictive maintenance",
        "",
        "• Telekomunikasi: Churn prediction dan customer retention"
    ])
    
    # Slide 7: Kasus Sederhana - Clustering
    add_content_slide(prs, "6. Kasus Sederhana: Data Mining (Clustering)", [
        "Studi Kasus: Segmentasi Pelanggan Toko Online",
        "",
        "Tujuan: Mengelompokkan pelanggan berdasarkan pola pembelian",
        "mereka untuk strategi pemasaran yang lebih efektif.",
        "",
        "Data yang digunakan:",
        "  • Jumlah pembelian per bulan",
        "  • Total pengeluaran per bulan",
        "  • Frekuensi kunjungan website"
    ])
    
    # Slide 8: Kasus - Metode Clustering
    add_content_slide(prs, "6. Kasus Sederhana: Metode K-Means Clustering", [
        "Algoritma K-Means:",
        "",
        "1. Tentukan jumlah cluster (K) yang diinginkan",
        "",
        "2. Inisialisasi K centroid secara random",
        "",
        "3. Hitung jarak setiap data ke centroid terdekat",
        "",
        "4. Kelompokkan data ke centroid terdekat",
        "",
        "5. Hitung ulang centroid dari rata-rata data dalam cluster",
        "",
        "6. Ulangi langkah 3-5 sampai centroid tidak berubah"
    ])
    
    # Slide 9: Kasus - Hasil Clustering
    add_content_slide(prs, "6. Kasus Sederhana: Hasil Clustering", [
        "Hasil Segmentasi Pelanggan:",
        "",
        "Cluster 1 - Pelanggan Premium:",
        "  Frekuensi pembelian tinggi, pengeluaran besar",
        "",
        "Cluster 2 - Pelanggan Regular:",
        "  Frekuensi pembelian sedang, pengeluaran sedang",
        "",
        "Cluster 3 - Pelanggan Casual:",
        "  Frekuensi pembelian rendah, pengeluaran kecil",
        "",
        "Strategi: Targeting berbeda untuk setiap cluster"
    ])
    
    # Slide 10: Kesimpulan
    add_content_slide(prs, "Kesimpulan", [
        "• Data Mining adalah alat penting untuk mengekstrak pengetahuan",
        "  dari data besar.",
        "",
        "• Memiliki berbagai fungsi seperti prediksi, klasifikasi, dan clustering.",
        "",
        "• Proses Data Mining melibatkan 7 tahapan utama dari definisi",
        "  masalah hingga deployment.",
        "",
        "• Memiliki banyak aplikasi praktis di berbagai industri.",
        "",
        "• Clustering adalah salah satu teknik penting untuk segmentasi data."
    ])
    
    # Simpan presentasi
    prs.save('Data_Mining.pptx')
    print("✅ Presentasi 'Data_Mining.pptx' berhasil dibuat!")
    print("📁 File tersimpan di direktori saat ini.")

if __name__ == "__main__":
    create_presentation()
